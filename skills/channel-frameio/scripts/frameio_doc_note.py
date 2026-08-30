# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=0.6", "openpyxl>=3.1"]
# ///
"""Scaffold a staged note for a Frame.io document capture (pdf/pptx/xlsx).

platform: frameio
scope: platform-general. Frame.io captures aren't HTML pages (no page.md —
just a downloaded document + meta.json), so the generic scaffold script
can't build these; this is the Frame.io equivalent of the plugin's
youtube_note.py.

Copies the document into the bundle's `assets/docs/` (docs are always
bundled — small), writes frontmatter + a TODO-SUMMARY placeholder + a
courtesy embed/link, and for pptx/xlsx dumps extracted text under a
collapsed "Extracted text" section so the agent writing the summary doesn't
need separate tooling. PDFs are left for the agent to Read directly (richer
than a naive text dump).

Usage:
  uv run frameio_doc_note.py <root> --capture-dir <dir> --notes-dir <dir> \
      [--group "<bundle name>"] [--group-type <kind>] [--author NAME]
      [--title-strip "<suffix>"] [--force]

`--title-strip` trims a share-wide suffix off every title (share viewers
often append the share's own name to each asset title); omit it to keep
titles as captured.

Prints {"note": <path relative to the wiki root>, "asset": <bundle-relative
path>}.

History:
  2026-07-15  created — batch speaker-slides share (25 docs, no page.md).
  2026-07-29  packaged into the channel-frameio skill unit: notes default to
              a domain-wide bundle root, the share-suffix trim is
              --title-strip, and provenance arrives at handoff complete
              (`generated`).
  2026-08-20  `--notes-dir` loses its default and becomes REQUIRED — the
              caller supplies the watch's own `dirs.sources` (the job's
              `dest`), never a domain-wide default a unit refresh could
              move.
"""

import argparse
import hashlib
import json
import re
import shutil
import sys
from pathlib import Path
from urllib.parse import urlsplit

SLUG_RE = re.compile(r"[^a-z0-9]+")
FILENAME_RE = re.compile(r"[^a-z0-9._-]+")


def source_hosts_for(host: str) -> list:
    """`www.a.com` -> `["www.a.com", "a.com"]` — progressively broader hosts so
    a query for either matches. No Public Suffix List; `co.uk` as a trailing
    entry is a correct host and a useless filter value, which is the cheaper
    trade. Duplicated from the pipeline's scaffold.py on purpose: a skill unit
    runs as a hosted script (the front door's `run` verb) and cannot import
    plugin/scripts/."""
    if not host:
        return []
    labels = host.split(".")
    return [".".join(labels[i:]) for i in range(max(1, len(labels) - 1))]


def slugify(parts, max_len=80):
    s = "-".join(parts).lower()
    s = SLUG_RE.sub("-", s).strip("-")
    return s[:max_len].rstrip("-") or "item"


def extract_pptx_text(path):
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        lines.append(t)
        slides.append(f"**Slide {i}**\n" + "\n".join(lines))
    return "\n\n".join(slides)


def extract_xlsx_text(path):
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    out = []
    for name in wb.sheetnames:
        ws = wb[name]
        out.append(f"**Sheet: {name}** ({ws.max_row} rows x {ws.max_column} cols)")
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 40), values_only=True):
            if any(c is not None for c in row):
                out.append(" | ".join("" if c is None else str(c) for c in row))
        if ws.max_row > 40:
            out.append(f"... ({ws.max_row - 40} more rows)")
    return "\n".join(out)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("root", type=Path, help="wiki root path")
    ap.add_argument("--capture-dir", required=True)
    ap.add_argument(
        "--notes-dir",
        default=None,
        help="wiki-relative bundle root. Defaults to the `dest` "
        "the host wrote onto capture.json — the watch's own "
        "dirs.sources. Pass it only for a hand-run capture "
        "with no job behind it",
    )
    ap.add_argument("--group", default=None)
    ap.add_argument("--group-type", dest="group_type", default=None)
    ap.add_argument("--author", default=None)
    ap.add_argument(
        "--title-strip", default=None, help="suffix to trim off captured titles (e.g. the share's own name)"
    )
    ap.add_argument("--force", action="store_true")
    args = ap.parse_args()

    cap_dir = args.root / args.capture_dir
    cap = json.loads((cap_dir / "capture.json").read_text())
    asset = cap["assets"][0]
    src_path = cap_dir / asset["local_path"]
    ext = src_path.suffix.lstrip(".").lower()
    path_bits = cap.get("path") or []
    title = (cap.get("title") or src_path.stem).strip()
    if args.title_strip:
        title = title.split(args.title_strip)[0].strip() or title

    # `--dest`-or-`capture.json` is `scaffold.py`'s rule, and this is the
    # same question: the HOST decides where a watch's notes land and writes
    # it onto the capture, so a unit reads it and never guesses. No third
    # fallback — a capture with neither is one `job_from_watch` never
    # produced, and a guessed bundle would silently disagree with wherever
    # the watch's notes actually go.
    notes_rel = args.notes_dir or cap.get("dest")
    if not notes_rel:
        sys.exit(
            f"{cap_dir / 'capture.json'} carries no `dest` and "
            f"--notes-dir was not given — every job-driven capture has "
            f"one (dirs.sources is required on every watch); pass "
            f"--notes-dir explicitly for a hand-run capture"
        )
    notes_dir = Path(notes_rel)
    domain_dir = args.root / notes_dir
    # path_bits[0] is the shared top folder every item in a batch carries —
    # including it wastes the slug's truncation budget on a prefix that
    # doesn't distinguish anything.
    slug = slugify(path_bits[1:] + [title])
    note_rel = notes_dir / "pages" / f"{slug}.md"
    note_path = args.root / note_rel
    if note_path.exists() and not args.force:
        sys.exit(f"{note_rel} exists — use --force")
    note_path.parent.mkdir(parents=True, exist_ok=True)

    # capture_asset.py names every doc generically ("document.<ext>"), so use
    # the real original filename (meta.json's "name") sanitized + content-
    # addressed by source URL — otherwise every PDF in the bundle collides on
    # "document.pdf" and silently dedupes onto the first one copied.
    meta = json.loads((cap_dir / "meta.json").read_text())
    orig_name = meta.get("name") or src_path.name
    stem = FILENAME_RE.sub("-", Path(orig_name).stem.lower()).strip("-") or "file"
    hash8 = hashlib.sha1((cap.get("url") or "").encode()).hexdigest()[:8]
    docs_dir = domain_dir / "assets" / "docs"
    docs_dir.mkdir(parents=True, exist_ok=True)
    dest = docs_dir / f"{hash8}-{stem}.{ext}"
    if not dest.exists():
        shutil.copy2(src_path, dest)
    asset_rel = f"../assets/docs/{dest.name}"

    extracted = ""
    if ext == "pptx":
        try:
            extracted = extract_pptx_text(src_path)
        except Exception as e:
            extracted = f"(extraction failed: {e})"
    elif ext == "xlsx":
        try:
            extracted = extract_xlsx_text(src_path)
        except Exception as e:
            extracted = f"(extraction failed: {e})"

    fm = ["---", f'title: "{title.replace(chr(34), chr(39))}"', f"resource: {cap.get('url')}", "type: doc"]
    # Scoping dimensions, each with exactly one meaning — the same set and the
    # same shapes `scaffold.py` emits, so a frameio note filters identically to
    # a generic scrape note.
    for key in ("author", "group", "group_type"):
        value = getattr(args, key, None) or cap.get(key) or ""
        if value:
            fm.append(f"{key}: {value}")
    # From the capture's own URL, not from the bundle directory. It used to
    # read `cap_dir.parent.name`, which was the netloc; that parent is the
    # watch's SLUG now, and `source_host` is a fact about the SOURCE — a
    # search filter — so it has to come from the URL either way.
    source_host = source_hosts_for(cap.get("domain") or urlsplit(cap.get("url") or "").netloc)
    areas = [f'"[[{a}]]"' for a in (cap.get("areas") or []) if a]
    fm.append(f"source_host: [{', '.join(source_host)}]")
    fm.append(f"areas: [{', '.join(areas)}]")
    fm.append(f"tags: [{', '.join(cap.get('tags') or [])}]")
    fm.append("status: draft")
    fm.append("---")

    breadcrumb = " / ".join(path_bits[1:]) if len(path_bits) > 1 else ""
    body = [f"# {title}"]
    if breadcrumb:
        body.append(f"*{breadcrumb}*")
    if ext == "pdf":
        body.append(f"![]({asset_rel})")
    else:
        body.append(f"[{orig_name}]({asset_rel})")
    if extracted:
        body.append("> [!note]- Extracted text\n" + "\n".join(f"> {line}" for line in extracted.splitlines()))

    placeholder = (
        "\n> [!summary]\n"
        "> TODO-SUMMARY — replace with 2-4 information-dense "
        "sentences before completing the handoff.\n\n"
    )
    note_path.write_text("\n".join(fm) + "\n" + placeholder + "\n\n".join(body) + "\n")

    print(
        json.dumps(
            {
                "note": str(note_rel),
                "asset": asset_rel,
                "ext": ext,
                "src_for_reading": str((cap_dir / asset["local_path"])),
            }
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
